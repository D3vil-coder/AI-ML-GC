"""
Chart Generator Agent (Agent 5)
Creates NATIVE PowerPoint charts (not matplotlib images).
Charts are fully editable in PowerPoint.
"""

import logging
from typing import Dict, List, Any, Tuple, Optional
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.brand_guidelines import BrandGuidelines

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ChartGenerator:
    """
    Generates NATIVE PowerPoint charts from financial data.
    All charts are fully editable (not images).
    """
    
    def __init__(self):
        self.brand = BrandGuidelines
    
    def create_revenue_ebitda_chart(self, slide, financials: Dict[str, Dict[int, float]],
                                     left: float, top: float, 
                                     width: float, height: float) -> Any:
        """
        Create Revenue & EBITDA combination chart.
        Revenue as bars, EBITDA as line.
        
        Args:
            slide: PowerPoint slide to add chart to
            financials: Financial data with 'revenue' and 'ebitda' keys
            left, top, width, height: Chart position and size in inches
            
        Returns:
            Chart shape object
        """
        revenue = financials.get('revenue', {})
        ebitda = financials.get('ebitda', {})
        
        if not revenue:
            logger.warning("No revenue data for chart")
            return None
        
        # Get years and values
        years = sorted(set(revenue.keys()) | set(ebitda.keys()))[-5:]  # Last 5 years
        
        categories = [f"FY{str(y)[-2:]}" for y in years]
        revenue_values = [revenue.get(y, 0) for y in years]
        ebitda_values = [ebitda.get(y, 0) for y in years]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Revenue (₹ Cr)', revenue_values)
        chart_data.add_series('EBITDA (₹ Cr)', ebitda_values)
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        # Style the chart
        self._style_chart(chart, "Revenue & EBITDA Trend")
        
        # Color the series
        plot = chart.plots[0]
        if len(plot.series) >= 1:
            plot.series[0].format.fill.solid()
            plot.series[0].format.fill.fore_color.rgb = self.brand.PRIMARY.rgb
        if len(plot.series) >= 2:
            plot.series[1].format.fill.solid()
            plot.series[1].format.fill.fore_color.rgb = self.brand.ACCENT.rgb
        
        # Add data labels
        for series in plot.series:
            series.has_data_labels = True
            series.data_labels.font.size = Pt(8)
            series.data_labels.font.color.rgb = self.brand.TEXT_DARK.rgb
        
        logger.info(f"Created Revenue & EBITDA chart with {len(years)} years of data")
        return chart
    
    def create_margin_donut_chart(self, slide, financials: Dict[str, Dict[int, float]],
                                   left: float, top: float,
                                   width: float, height: float) -> Any:
        """
        Create margin donut chart showing profitability metrics.
        
        Args:
            slide: PowerPoint slide
            financials: Financial data
            left, top, width, height: Position and size
            
        Returns:
            Chart shape
        """
        # Get latest year data
        pat_margin = financials.get('pat_margin', {})
        ebitda = financials.get('ebitda', {})
        revenue = financials.get('revenue', {})
        
        # Calculate margins for latest year
        common_years = set(ebitda.keys()) & set(revenue.keys())
        if not common_years and not pat_margin:
            logger.warning("Insufficient data for margin chart")
            return None
        
        margins = []
        labels = []
        
        if common_years:
            latest = max(common_years)
            if revenue[latest] > 0:
                ebitda_margin = (ebitda[latest] / revenue[latest]) * 100
                margins.append(ebitda_margin)
                labels.append('EBITDA %')
                margins.append(100 - ebitda_margin)
                labels.append('Other Costs')
        
        if not margins:
            # Default visualization
            margins = [20, 80]
            labels = ['Margin', 'Costs']
        
        # Create doughnut chart
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series('Margin Split', margins)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        self._style_chart(chart, "Margin Analysis")
        
        # Color the slices
        plot = chart.plots[0]
        series = plot.series[0]
        colors = [self.brand.PRIMARY.rgb, RGBColor(200, 200, 200)]
        
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]
        
        logger.info("Created margin donut chart")
        return chart
    
    def create_line_chart(self, slide, data_series: Dict[str, Dict[int, float]],
                          title: str,
                          left: float, top: float,
                          width: float, height: float) -> Any:
        """
        Create a line chart for trend visualization.
        
        Args:
            slide: PowerPoint slide
            data_series: Dict of series name -> {year: value}
            title: Chart title
            left, top, width, height: Position and size
            
        Returns:
            Chart shape
        """
        if not data_series:
            return None
        
        # Get all years across all series
        all_years = set()
        for series_data in data_series.values():
            all_years.update(series_data.keys())
        years = sorted(all_years)[-5:]  # Last 5 years
        
        categories = [f"FY{str(y)[-2:]}" for y in years]
        
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for series_name, values in data_series.items():
            series_values = [values.get(y, 0) for y in years]
            chart_data.add_series(series_name, series_values)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        self._style_chart(chart, title)
        
        # Color series
        colors = [self.brand.PRIMARY.rgb, self.brand.SECONDARY.rgb, self.brand.ACCENT.rgb]
        for i, series in enumerate(chart.plots[0].series):
            series.format.line.color.rgb = colors[i % len(colors)]
            series.format.line.width = Pt(2)
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = colors[i % len(colors)]
        
        logger.info(f"Created line chart: {title}")
        return chart
    
    def create_kpi_table(self, slide, kpis: Dict[str, str],
                         left: float, top: float,
                         width: float, height: float) -> Any:
        """
        Create a KPI table as a text box styled table.
        
        Args:
            slide: PowerPoint slide
            kpis: Dict of KPI name -> value
            left, top, width, height: Position and size
            
        Returns:
            Shape object
        """
        if not kpis:
            return None
        
        # Create table
        rows = len(kpis)
        cols = 2
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top), Inches(width), Inches(height)
        ).table
        
        # Style and populate
        row_idx = 0
        for kpi_name, kpi_value in kpis.items():
            table.cell(row_idx, 0).text = kpi_name
            table.cell(row_idx, 1).text = str(kpi_value)
            
            # Style cells
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
                
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.name = self.brand.FONT_BODY
                        run.font.color.rgb = self.brand.TEXT_DARK.rgb
            
            row_idx += 1
        
        logger.info(f"Created KPI table with {rows} metrics")
        return table
    
    def _style_chart(self, chart, title: str):
        """Apply standard styling to chart."""
        # Title
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.name = self.brand.FONT_HEADING
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.brand.TEXT_DARK.rgb
        
        # Legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = self.brand.FONT_BODY
    
    def calculate_cagr(self, values: Dict[int, float]) -> Optional[float]:
        """Calculate Compound Annual Growth Rate."""
        if len(values) < 2:
            return None
        
        years = sorted(values.keys())
        first_year = years[0]
        last_year = years[-1]
        
        first_value = values[first_year]
        last_value = values[last_year]
        
        if first_value <= 0:
            return None
        
        n_years = last_year - first_year
        if n_years <= 0:
            return None
        
        cagr = ((last_value / first_value) ** (1 / n_years) - 1) * 100
        return cagr


# Test
if __name__ == "__main__":
    from pptx import Presentation
    
    # Create test presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
    
    # Test data
    financials = {
        'revenue': {2020: 100, 2021: 120, 2022: 145, 2023: 175, 2024: 210},
        'ebitda': {2020: 15, 2021: 20, 2022: 25, 2023: 32, 2024: 40},
        'pat_margin': {2024: 8.5},
    }
    
    generator = ChartGenerator()
    
    # Create charts
    generator.create_revenue_ebitda_chart(slide, financials, 0.5, 1.0, 4.5, 3.5)
    generator.create_margin_donut_chart(slide, financials, 5.5, 1.0, 4.0, 3.5)
    
    # Create KPI table
    kpis = {
        'Revenue CAGR': '20.4%',
        'EBITDA Margin': '19.0%',
        'ROE': '15.2%',
        'ROCE': '18.5%',
    }
    generator.create_kpi_table(slide, kpis, 0.5, 5.0, 9.0, 1.5)
    
    # Save
    output_path = "test_charts.pptx"
    prs.save(output_path)
    print(f"Saved test presentation to: {output_path}")
    print("Open the file and verify charts are editable (right-click -> Edit Data)")
